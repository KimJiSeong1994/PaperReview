"""
Poster Exporter

HTML 포스터를 다양한 형식(PDF, PPTX)으로 내보내기
Paper2Poster의 출력 형식 다양화 구현
"""

from pathlib import Path
from typing import Optional
import subprocess
import os


class PosterExporter:
    """
    포스터 내보내기 유틸리티
    
    지원 형식:
    - HTML (기본)
    - PDF (wkhtmltopdf 또는 playwright 사용)
    - PPTX (python-pptx 사용, 선택)
    """
    
    def __init__(self):
        self.has_playwright = self._check_playwright()
        self.has_pptx = self._check_pptx()
    
    def _check_playwright(self) -> bool:
        """Playwright 사용 가능 여부 확인"""
        try:
            from playwright.sync_api import sync_playwright
            return True
        except ImportError:
            return False
    
    def _check_pptx(self) -> bool:
        """python-pptx 사용 가능 여부 확인"""
        try:
            from pptx import Presentation  # type: ignore
            return True
        except ImportError:
            return False
    
    def export_to_pdf(self, html_path: Path, output_path: Optional[Path] = None) -> Optional[Path]:
        """
        HTML을 PDF로 변환
        
        Args:
            html_path: 입력 HTML 파일 경로
            output_path: 출력 PDF 경로 (기본값: html_path와 동일 경로에 .pdf)
            
        Returns:
            생성된 PDF 파일 경로
        """
        if output_path is None:
            output_path = html_path.with_suffix('.pdf')
        
        # Method 1: Playwright (권장)
        if self.has_playwright:
            try:
                return self._export_pdf_playwright(html_path, output_path)
            except Exception as e:
                pass
        
        # Method 2: wkhtmltopdf (fallback)
        try:
            return self._export_pdf_wkhtmltopdf(html_path, output_path)
        except Exception as e:
            pass
        
        return None
    
    def _export_pdf_playwright(self, html_path: Path, output_path: Path) -> Path:
        """Playwright를 사용한 PDF 변환"""
        from playwright.sync_api import sync_playwright
        
        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page()
            
            # HTML 로드
            page.goto(f"file://{html_path.absolute()}")
            
            # PDF 생성 (landscape, A3 사이즈)
            page.pdf(
                path=str(output_path),
                format='A3',
                landscape=True,
                print_background=True,
                margin={
                    'top': '0',
                    'right': '0',
                    'bottom': '0',
                    'left': '0'
                }
            )
            
            browser.close()
        
        return output_path
    
    def _export_pdf_wkhtmltopdf(self, html_path: Path, output_path: Path) -> Path:
        """wkhtmltopdf를 사용한 PDF 변환"""
        cmd = [
            'wkhtmltopdf',
            '--page-size', 'A3',
            '--orientation', 'Landscape',
            '--enable-local-file-access',
            str(html_path),
            str(output_path)
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode != 0:
            raise RuntimeError(f"wkhtmltopdf failed: {result.stderr}")
        
        return output_path
    
    def export_to_pptx(self, html_content: str, output_path: Path, title: str = "Research Poster") -> Optional[Path]:
        """
        HTML 콘텐츠를 PPTX로 변환
        
        Args:
            html_content: HTML 문자열
            output_path: 출력 PPTX 경로
            title: 슬라이드 제목
            
        Returns:
            생성된 PPTX 파일 경로
        """
        if not self.has_pptx:
            return None
        
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches, Pt  # type: ignore
            from pptx.enum.text import PP_ALIGN  # type: ignore
            
            # 프레젠테이션 생성 (16:9 비율)
            prs = Presentation()
            prs.slide_width = Inches(13.33)  # Wide format
            prs.slide_height = Inches(7.5)
            
            # 빈 슬라이드 추가
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # 제목 추가
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(12.33)
            height = Inches(1)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            
            # 제목 스타일
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # 내용 추가 (간단한 텍스트 추출)
            content_text = self._extract_text_from_html(html_content)
            
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(12.33)
            height = Inches(5.5)
            
            content_box = slide.shapes.add_textbox(left, top, width, height)
            content_frame = content_box.text_frame
            content_frame.text = content_text[:1000]  # 제한된 길이
            content_frame.word_wrap = True
            
            # 저장
            prs.save(str(output_path))
            
            return output_path
            
        except Exception as e:
            return None
    
    def _extract_text_from_html(self, html_content: str) -> str:
        """HTML에서 텍스트 추출 (간단한 파싱)"""
        import re
        
        # 태그 제거
        text = re.sub(r'<style[^>]*>.*?</style>', '', html_content, flags=re.DOTALL)
        text = re.sub(r'<script[^>]*>.*?</script>', '', text, flags=re.DOTALL)
        text = re.sub(r'<[^>]+>', ' ', text)
        
        # 공백 정리
        text = re.sub(r'\s+', ' ', text)
        
        return text.strip()
    
    def export_all_formats(self, html_path: Path, base_name: Optional[str] = None) -> dict:
        """
        모든 형식으로 동시 내보내기
        
        Args:
            html_path: HTML 파일 경로
            base_name: 기본 파일명 (없으면 html_path 사용)
            
        Returns:
            dict: {'pdf': Path, 'pptx': Path}
        """
        if base_name is None:
            base_name = html_path.stem
        
        output_dir = html_path.parent
        
        results = {
            'html': html_path,
            'pdf': None,
            'pptx': None
        }
        
        # PDF 생성
        pdf_path = output_dir / f"{base_name}.pdf"
        results['pdf'] = self.export_to_pdf(html_path, pdf_path)
        
        # PPTX 생성 (HTML 읽기 필요)
        if self.has_pptx:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            pptx_path = output_dir / f"{base_name}.pptx"
            results['pptx'] = self.export_to_pptx(html_content, pptx_path, base_name)
        
        return results


# 전역 익스포터 인스턴스
_exporter = None

def get_exporter() -> PosterExporter:
    """싱글톤 익스포터 인스턴스 반환"""
    global _exporter
    if _exporter is None:
        _exporter = PosterExporter()
    return _exporter

